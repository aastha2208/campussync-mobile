from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


OUT_DIR = Path(__file__).resolve().parent
TEAM = "Tanishk, Aastha, Aditi"
TOPIC = "Secure Data Transmission over Netcat: Demonstrating Confidentiality through Plaintext vs Encrypted Channels"


REPORT_MD = f"""# {TOPIC}

**Assignment Title:** Practical Exploration of Cryptography / Security Tools  
**Selected Tool:** Netcat / Ncat  
**Team Members:** {TEAM}

---

## Abstract

This mini project demonstrates why confidential data should not be sent using a normal plaintext network channel. Netcat is a simple networking tool that can send messages or files between two systems. In its normal form, Netcat sends data as plaintext, which means an attacker using a packet capture tool such as Wireshark can read the transmitted message. To show the secure alternative, we use Ncat, a Netcat-compatible tool that supports SSL/TLS encryption. The same message is transmitted once through a plaintext Netcat channel and once through an encrypted Ncat SSL channel. The captured traffic is then compared. The result clearly shows that plaintext traffic is readable, while encrypted traffic is not understandable to the attacker. This proves the importance of confidentiality in secure data transmission.

---

## 1. Introduction of the Tool

Netcat is a command-line networking utility used to read from and write to network connections using TCP or UDP. It is often called the "Swiss Army knife" of networking because it can be used for testing ports, sending messages, transferring files, and debugging network services.

For this project, Netcat is used to create a simple sender-receiver communication channel. The normal Netcat command sends data in plaintext. This is useful for learning, but unsafe for confidential information. To demonstrate secure transmission, we use Ncat with SSL support. Ncat is part of the Nmap project and is compatible with Netcat-style usage, but it also supports encrypted communication using SSL/TLS.

In simple words:

- **Netcat (`nc`)**: sends data normally, but without encryption.
- **Ncat (`ncat --ssl`)**: sends data through an encrypted SSL/TLS channel.
- **Wireshark**: captures packets and helps us compare plaintext and encrypted traffic.

---

## 2. Objective of the Project

The main objective is to demonstrate confidentiality during data transmission.

Specific objectives:

1. To understand how Netcat sends data between two systems.
2. To demonstrate that plaintext Netcat traffic can be read by an attacker.
3. To demonstrate encrypted data transmission using Ncat SSL.
4. To compare plaintext traffic and encrypted traffic using Wireshark.
5. To explain why encryption is important in cryptography and network security.

---

## 3. Installation and Setup

### Required Tools

| Tool | Purpose |
|---|---|
| Netcat / Ncat | To send data between sender and receiver |
| Wireshark | To capture and inspect network packets |
| OpenSSL | To generate a temporary SSL certificate for encrypted demo |
| Two terminals | One terminal acts as receiver, one acts as sender |

### Recommended Setup

The easiest setup is one laptop with two terminals:

- Terminal 1: Receiver
- Terminal 2: Sender
- Wireshark: Packet capture

If two laptops are available, one laptop can be sender and the other can be receiver. Both should be connected to the same Wi-Fi network.

### Ubuntu / Linux Installation Commands

```bash
sudo apt update
sudo apt install netcat-openbsd ncat wireshark openssl -y
```

### Windows Installation

1. Install Nmap from: https://nmap.org/download.html  
   Ncat is included with Nmap.
2. Install Wireshark from: https://www.wireshark.org/download.html
3. During Wireshark/Npcap installation, enable loopback capture if using one laptop.

---

## 4. Basic Theory

### 4.1 What is Plaintext?

Plaintext means readable data. If we send the message "CRP secret data" without encryption, anyone who captures the packet can read the exact message.

### 4.2 What is Encryption?

Encryption converts readable data into unreadable ciphertext using cryptographic algorithms. Only the correct receiver can decrypt it.

### 4.3 What is Confidentiality?

Confidentiality is one of the main goals of cybersecurity. It means that data should be visible only to authorized users. In this project, confidentiality is shown by comparing readable plaintext traffic with unreadable encrypted traffic.

### 4.4 Why Netcat Alone is Not Secure

Traditional Netcat does not provide encryption. It simply sends bytes over the network. Therefore, it is useful for testing but should not be used alone to send passwords, private messages, or sensitive files.

### 4.5 Why Ncat SSL is Used

Ncat is a modern Netcat-family tool. It supports SSL/TLS encryption. TLS is the same basic security technology used by HTTPS websites. In this project, Ncat SSL is used to show how encryption protects data in transit.

---

## 5. Mini Project Architecture

The project has three roles:

1. **Sender**: Sends a secret message/file.
2. **Receiver**: Receives the message/file.
3. **Attacker/Observer**: Uses Wireshark to capture packets and check whether the message is readable.

### Plaintext Flow

Sender -> Plaintext Netcat Channel -> Receiver  
Attacker using Wireshark -> Can read the message

### Encrypted Flow

Sender -> Encrypted Ncat SSL Channel -> Receiver  
Attacker using Wireshark -> Cannot read the message

---

## 6. Practical Demonstration

### Important Safety Note

This project is for educational demonstration only. Run it on your own laptop, lab computers, or with permission on a local network. Do not capture other people's traffic.

### Demo Message

Create a file named `secret_message.txt`:

```text
This is confidential CRP data from Tanishk, Aastha and Aditi.
If this line is visible in Wireshark, confidentiality is broken.
```

---

## 7. Experiment 1: Plaintext Data Transmission using Netcat

### Step 1: Start Wireshark Capture

Open Wireshark and start capturing on:

- `lo` / loopback interface if using one Linux laptop with `127.0.0.1`
- Wi-Fi interface if using two laptops

Use this display filter:

```text
tcp.port == 4444
```

### Step 2: Start Receiver

In Terminal 1:

```bash
nc -lvnp 4444 > received_plain.txt
```

### Step 3: Send File from Sender

In Terminal 2:

```bash
nc 127.0.0.1 4444 < secret_message.txt
```

If using two laptops, replace `127.0.0.1` with the receiver laptop IP address.

### Step 4: Check Received File

```bash
cat received_plain.txt
```

### Step 5: Wireshark Observation

In Wireshark:

1. Apply filter: `tcp.port == 4444`
2. Right-click a packet.
3. Select **Follow > TCP Stream**.
4. Observe that the secret message is visible in readable form.

### Result

The confidential message is visible in Wireshark. This proves that normal Netcat does not provide confidentiality.

---

## 8. Cryptanalysis Demonstration

In this project, cryptanalysis is demonstrated as traffic inspection/packet sniffing.

In the plaintext experiment, the attacker does not need to break any encryption because there is no encryption. The attacker simply captures packets using Wireshark and follows the TCP stream. The secret message appears directly on screen.

This is a practical attack because many real-world attackers use packet sniffing on unsafe networks to steal plaintext information.

### Cryptanalysis Observation for Plaintext Channel

| Item | Observation |
|---|---|
| Capture tool | Wireshark |
| Port observed | 4444 |
| Message readable? | Yes |
| Confidentiality protected? | No |
| Attack difficulty | Easy |

---

## 9. Experiment 2: Encrypted Transmission using Ncat SSL

### Step 1: Generate Temporary SSL Certificate

Run this once:

```bash
openssl req -x509 -newkey rsa:2048 -nodes -keyout server.key -out server.crt -days 1 -subj "/CN=NetcatAAT"
```

This creates:

- `server.key`: private key
- `server.crt`: certificate

For this class demo, a temporary self-signed certificate is enough.

### Step 2: Start Wireshark Capture

Use this display filter:

```text
tcp.port == 4445
```

### Step 3: Start Encrypted Receiver

In Terminal 1:

```bash
ncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445 > received_secure.txt
```

### Step 4: Send File through Encrypted Channel

In Terminal 2:

```bash
ncat --ssl 127.0.0.1 4445 < secret_message.txt
```

If using two laptops, replace `127.0.0.1` with the receiver laptop IP address.

### Step 5: Check Received File

```bash
cat received_secure.txt
```

### Step 6: Wireshark Observation

In Wireshark:

1. Apply filter: `tcp.port == 4445`
2. Right-click a packet.
3. Select **Follow > TCP Stream**.
4. Observe that the secret message is not visible in readable form.

### Result

The receiver gets the correct message, but the captured packets do not show readable data. This proves that encryption provides confidentiality.

---

## 10. Comparison Table

| Feature | Plaintext Netcat | Encrypted Ncat SSL |
|---|---|---|
| Command used | `nc` | `ncat --ssl` |
| Port | 4444 | 4445 |
| Data visible in Wireshark | Yes | No |
| Confidentiality | Not provided | Provided |
| Attacker can read message | Yes | No |
| Suitable for sensitive data | No | Better than plaintext |
| Learning value | Shows insecurity | Shows protection |

---

## 11. Application and Relevance to Cryptography / Security

This project is relevant because many applications send data over networks. If the data is not encrypted, attackers may capture and read it. Encryption protects confidentiality.

Real-world applications:

1. HTTPS protects website data.
2. SSH protects remote login sessions.
3. VPNs protect network traffic.
4. Secure file transfer tools protect files during transmission.
5. Messaging applications use encryption to protect private chats.

This project shows the same concept in a simple and practical way using Netcat-style commands.

---

## 12. Limitations

1. Traditional Netcat does not support encryption by default.
2. Ncat SSL is useful for demonstration but not a full replacement for secure tools like SSH or SCP.
3. Self-signed certificates are acceptable for lab demos but not ideal for production systems.
4. The project focuses mainly on confidentiality, not authentication or integrity in full detail.

---

## 13. Team Contribution

| Member | Suggested Responsibility |
|---|---|
| Tanishk | Explain Netcat basics and perform plaintext demo |
| Aastha | Explain confidentiality, encryption, and Ncat SSL demo |
| Aditi | Explain Wireshark observations, comparison table, conclusion |

---

## 14. Conclusion

The project successfully demonstrates secure data transmission over a Netcat-style tool. In the first experiment, normal Netcat sends data as plaintext, and the secret message is easily visible in Wireshark. In the second experiment, Ncat SSL encrypts the channel, so the same message is delivered to the receiver but is not readable in packet capture. Therefore, encryption protects confidentiality during data transmission.

Final conclusion: **Plaintext channels are unsafe for confidential data, while encrypted channels protect data from packet sniffing attacks.**

---

## 15. References

1. Netcat manual pages
2. Nmap Ncat documentation: https://nmap.org/ncat/
3. Wireshark documentation: https://www.wireshark.org/docs/
4. OpenSSL documentation: https://www.openssl.org/docs/

---

## Appendix A: Screenshot Checklist

Add these screenshots in the final report/PPT after running the demo:

1. Terminal showing plaintext receiver command.
2. Terminal showing plaintext sender command.
3. Wireshark TCP stream showing readable secret message on port 4444.
4. Terminal showing certificate generation.
5. Terminal showing encrypted receiver command.
6. Terminal showing encrypted sender command.
7. Wireshark TCP stream showing unreadable encrypted data on port 4445.
8. Comparison table.
"""


DEMO_GUIDE = f"""# Beginner-Friendly Demo Guide

## Project

**{TOPIC}**

Team: {TEAM}

This guide is written for beginners. Follow it step by step.

---

## What You Need to Prove

You only need to prove one main idea:

> Normal Netcat sends readable data. Encrypted Ncat sends unreadable data to an attacker.

This directly satisfies the AAT requirement:

- Introduction of tool
- Installation/setup
- Practical demonstration
- Cryptanalysis demonstration
- Application/relevance
- Conclusion
- Report preparation

---

## Best Setup for Finishing Today

Use **one laptop** with:

- Terminal 1 as receiver
- Terminal 2 as sender
- Wireshark as attacker/observer

Use IP address `127.0.0.1`.

If Wireshark does not show loopback packets on Windows, use two laptops on the same Wi-Fi instead.

---

## Create the Secret File

```bash
echo "This is confidential CRP data from Tanishk, Aastha and Aditi." > secret_message.txt
echo "If this line is visible in Wireshark, confidentiality is broken." >> secret_message.txt
```

---

## Demo 1: Plaintext Netcat

### Terminal 1: Receiver

```bash
nc -lvnp 4444 > received_plain.txt
```

### Terminal 2: Sender

```bash
nc 127.0.0.1 4444 < secret_message.txt
```

### Wireshark

Filter:

```text
tcp.port == 4444
```

Right-click a packet -> Follow -> TCP Stream.

### What to Say

\"Here we used normal Netcat. The receiver got the message, but Wireshark also shows the same message clearly. So confidentiality is not protected.\"

---

## Demo 2: Encrypted Ncat SSL

### Generate Certificate

```bash
openssl req -x509 -newkey rsa:2048 -nodes -keyout server.key -out server.crt -days 1 -subj "/CN=NetcatAAT"
```

### Terminal 1: Encrypted Receiver

```bash
ncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445 > received_secure.txt
```

### Terminal 2: Encrypted Sender

```bash
ncat --ssl 127.0.0.1 4445 < secret_message.txt
```

### Wireshark

Filter:

```text
tcp.port == 4445
```

Right-click a packet -> Follow -> TCP Stream.

### What to Say

\"Now we sent the same file using Ncat SSL. The receiver still receives the correct message, but Wireshark cannot show the readable text. This proves confidentiality using encryption.\"

---

## Screenshot Checklist

Take these screenshots:

1. Plaintext receiver terminal.
2. Plaintext sender terminal.
3. Wireshark plaintext TCP stream showing readable message.
4. Certificate generation command.
5. Encrypted receiver terminal.
6. Encrypted sender terminal.
7. Wireshark encrypted TCP stream showing unreadable data.
8. Final comparison table from PPT/report.

---

## 3-Person Presentation Split

### Tanishk

- Introduce topic and Netcat.
- Explain what plaintext means.
- Run/show Demo 1.

### Aastha

- Explain confidentiality and encryption.
- Explain why Ncat SSL is used.
- Run/show Demo 2.

### Aditi

- Explain Wireshark observation.
- Explain cryptanalysis and comparison table.
- Give conclusion.

---

## Simple Viva Answers

**Q1. What is Netcat?**  
Netcat is a command-line networking tool used to send and receive data over TCP or UDP connections.

**Q2. What is the security problem with normal Netcat?**  
Normal Netcat sends data in plaintext, so captured traffic can be read.

**Q3. What is confidentiality?**  
Confidentiality means only authorized people should be able to read the data.

**Q4. What did Wireshark prove?**  
Wireshark proved that plaintext Netcat traffic is readable, while encrypted Ncat SSL traffic is not readable.

**Q5. Why did you use Ncat SSL?**  
Traditional Netcat does not provide encryption. Ncat is a Netcat-compatible tool that supports SSL/TLS encryption.

**Q6. What is the conclusion?**  
Plaintext transmission is unsafe for confidential data. Encrypted transmission protects the message from packet sniffing.
"""


VIVA_QA = """# Viva / Oral Presentation Questions and Answers

## 1. What is the main aim of your project?

The aim is to demonstrate confidentiality in data transmission by comparing plaintext Netcat communication with encrypted Ncat SSL communication.

## 2. What is Netcat?

Netcat is a command-line tool used to read and write data over network connections. It can be used for simple client-server communication, file transfer, and network testing.

## 3. Why is normal Netcat not secure?

Normal Netcat does not encrypt data. If someone captures the packets, they can read the message directly.

## 4. What is plaintext?

Plaintext is readable data before encryption. Example: a message that can be read exactly as typed.

## 5. What is ciphertext?

Ciphertext is encrypted data that looks unreadable to unauthorized people.

## 6. What is confidentiality?

Confidentiality means protecting information so that only authorized people can read it.

## 7. Which tool did you use for cryptanalysis?

We used Wireshark for packet capture and traffic inspection.

## 8. What did the plaintext experiment prove?

It proved that data sent using normal Netcat can be captured and read in Wireshark.

## 9. What did the encrypted experiment prove?

It proved that Ncat SSL hides the message from packet capture, while still delivering the correct file to the receiver.

## 10. Why did you use Ncat?

Ncat is a Netcat-compatible tool from the Nmap project. It supports SSL/TLS encryption, which traditional Netcat does not provide.

## 11. Is this safe for real production use?

It is good for learning and demonstration. For real systems, tools like SSH, SCP, HTTPS, or VPNs are preferred.

## 12. What is your final conclusion?

Plaintext channels are unsafe for confidential data. Encrypted channels protect data from packet sniffing and provide confidentiality.
"""


def add_doc_heading(document, text, level=1):
    document.add_heading(text, level=level)


def add_doc_paragraph(document, text="", bold=False):
    p = document.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    return p


def add_doc_bullets(document, items):
    for item in items:
        document.add_paragraph(item, style="List Bullet")


def add_code_block(document, code):
    for line in code.strip("\n").splitlines():
        p = document.add_paragraph()
        run = p.add_run(line)
        run.font.name = "Courier New"
        run.font.size = Pt(9)


def create_report_docx():
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("ALTERNATE ASSESSMENT TOOL REPORT")
    run.bold = True
    run.font.size = Pt(16)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(TOPIC)
    run.bold = True
    run.font.size = Pt(14)

    info = doc.add_paragraph()
    info.alignment = WD_ALIGN_PARAGRAPH.CENTER
    info.add_run("Assignment Title: ").bold = True
    info.add_run("Practical Exploration of Cryptography / Security Tools\n")
    info.add_run("Selected Tool: ").bold = True
    info.add_run("Netcat / Ncat\n")
    info.add_run("Team Members: ").bold = True
    info.add_run(TEAM)

    add_doc_heading(doc, "Abstract", 1)
    add_doc_paragraph(
        doc,
        "This mini project demonstrates why confidential data should not be sent using a normal plaintext network channel. "
        "Netcat is used to send data between a sender and receiver. In normal mode, the message is visible in Wireshark. "
        "Then the same message is sent using Ncat SSL, a Netcat-compatible encrypted channel. The receiver gets the data, "
        "but Wireshark cannot show the readable message. This proves confidentiality through encryption.",
    )

    add_doc_heading(doc, "1. Introduction of the Tool", 1)
    add_doc_paragraph(
        doc,
        "Netcat is a command-line networking utility used to read from and write to network connections using TCP or UDP. "
        "It can be used for testing ports, sending messages, transferring files, and debugging simple network services.",
    )
    add_doc_paragraph(
        doc,
        "Traditional Netcat does not encrypt data. To demonstrate an encrypted channel, this project uses Ncat with SSL. "
        "Ncat is a Netcat-compatible tool from the Nmap project and supports SSL/TLS encryption.",
    )
    add_doc_bullets(
        doc,
        [
            "Netcat (nc): sends data normally without encryption.",
            "Ncat SSL: sends data through an encrypted SSL/TLS channel.",
            "Wireshark: captures packets to check whether the message is readable.",
        ],
    )

    add_doc_heading(doc, "2. Objective", 1)
    add_doc_bullets(
        doc,
        [
            "Understand how Netcat sends data between two systems.",
            "Show that plaintext Netcat traffic can be read using Wireshark.",
            "Show encrypted data transmission using Ncat SSL.",
            "Compare plaintext and encrypted traffic.",
            "Explain the importance of confidentiality in cryptography/security.",
        ],
    )

    add_doc_heading(doc, "3. Installation and Setup", 1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "Tool"
    hdr[1].text = "Purpose"
    for tool, purpose in [
        ("Netcat / Ncat", "Send data between sender and receiver"),
        ("Wireshark", "Capture and inspect packets"),
        ("OpenSSL", "Generate a temporary SSL certificate"),
        ("Two terminals", "Run sender and receiver commands"),
    ]:
        row = table.add_row().cells
        row[0].text = tool
        row[1].text = purpose
    add_doc_paragraph(doc, "Ubuntu/Linux installation command:")
    add_code_block(doc, "sudo apt update\nsudo apt install netcat-openbsd ncat wireshark openssl -y")

    add_doc_heading(doc, "4. Basic Theory", 1)
    add_doc_paragraph(doc, "Plaintext means readable data. Ciphertext means encrypted unreadable data.")
    add_doc_paragraph(
        doc,
        "Confidentiality means only authorized users should be able to read the data. Encryption is used to achieve confidentiality.",
    )
    add_doc_paragraph(
        doc,
        "In this project, confidentiality is tested by capturing both channels in Wireshark and checking whether the secret message is visible.",
    )

    add_doc_heading(doc, "5. Mini Project Architecture", 1)
    add_doc_bullets(
        doc,
        [
            "Sender: sends a secret message or file.",
            "Receiver: receives the message or file.",
            "Attacker/Observer: uses Wireshark to capture packets.",
        ],
    )
    add_doc_paragraph(doc, "Plaintext flow: Sender -> Netcat plaintext channel -> Receiver. Wireshark can read the message.")
    add_doc_paragraph(doc, "Encrypted flow: Sender -> Ncat SSL encrypted channel -> Receiver. Wireshark cannot read the message.")

    add_doc_heading(doc, "6. Practical Demonstration", 1)
    add_doc_paragraph(doc, "Create the secret file:")
    add_code_block(
        doc,
        'echo "This is confidential CRP data from Tanishk, Aastha and Aditi." > secret_message.txt\n'
        'echo "If this line is visible in Wireshark, confidentiality is broken." >> secret_message.txt',
    )

    add_doc_heading(doc, "7. Experiment 1: Plaintext Netcat", 1)
    add_doc_paragraph(doc, "Receiver command:")
    add_code_block(doc, "nc -lvnp 4444 > received_plain.txt")
    add_doc_paragraph(doc, "Sender command:")
    add_code_block(doc, "nc 127.0.0.1 4444 < secret_message.txt")
    add_doc_paragraph(
        doc,
        "In Wireshark, use filter tcp.port == 4444, then Follow TCP Stream. The secret message is visible.",
    )

    add_doc_heading(doc, "8. Cryptanalysis Demonstration", 1)
    add_doc_paragraph(
        doc,
        "Here, cryptanalysis is shown as packet sniffing. In plaintext mode there is no encryption to break. "
        "The attacker captures packets using Wireshark and directly reads the message. This shows that confidentiality is broken.",
    )

    add_doc_heading(doc, "9. Experiment 2: Encrypted Ncat SSL", 1)
    add_doc_paragraph(doc, "Generate temporary certificate:")
    add_code_block(doc, 'openssl req -x509 -newkey rsa:2048 -nodes -keyout server.key -out server.crt -days 1 -subj "/CN=NetcatAAT"')
    add_doc_paragraph(doc, "Encrypted receiver command:")
    add_code_block(doc, "ncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445 > received_secure.txt")
    add_doc_paragraph(doc, "Encrypted sender command:")
    add_code_block(doc, "ncat --ssl 127.0.0.1 4445 < secret_message.txt")
    add_doc_paragraph(
        doc,
        "In Wireshark, use filter tcp.port == 4445, then Follow TCP Stream. The original message is not visible.",
    )

    add_doc_heading(doc, "10. Observation and Comparison", 1)
    comp = doc.add_table(rows=1, cols=3)
    comp.style = "Table Grid"
    hdr = comp.rows[0].cells
    hdr[0].text = "Feature"
    hdr[1].text = "Plaintext Netcat"
    hdr[2].text = "Encrypted Ncat SSL"
    rows = [
        ("Command", "nc", "ncat --ssl"),
        ("Port", "4444", "4445"),
        ("Visible in Wireshark", "Yes", "No"),
        ("Confidentiality", "Not provided", "Provided"),
        ("Attacker can read message", "Yes", "No"),
        ("Suitable for sensitive data", "No", "Better than plaintext"),
    ]
    for row_values in rows:
        row = comp.add_row().cells
        for i, value in enumerate(row_values):
            row[i].text = value

    add_doc_heading(doc, "11. Application and Relevance", 1)
    add_doc_bullets(
        doc,
        [
            "HTTPS protects website communication.",
            "SSH protects remote login sessions.",
            "VPNs protect network traffic on unsafe networks.",
            "Secure file transfer tools protect files during transmission.",
            "Encrypted messaging applications protect private chats.",
        ],
    )

    add_doc_heading(doc, "12. Limitations", 1)
    add_doc_bullets(
        doc,
        [
            "Traditional Netcat does not support encryption by default.",
            "Ncat SSL is useful for demonstration, but production systems usually use SSH, SCP, HTTPS, or VPN.",
            "Self-signed certificates are acceptable for lab demos but not ideal for real public systems.",
            "The project focuses mainly on confidentiality.",
        ],
    )

    add_doc_heading(doc, "13. Team Contribution", 1)
    team_table = doc.add_table(rows=1, cols=2)
    team_table.style = "Table Grid"
    team_table.rows[0].cells[0].text = "Member"
    team_table.rows[0].cells[1].text = "Responsibility"
    for member, role in [
        ("Tanishk", "Netcat introduction and plaintext demo"),
        ("Aastha", "Confidentiality concept and encrypted Ncat SSL demo"),
        ("Aditi", "Wireshark observation, comparison, and conclusion"),
    ]:
        row = team_table.add_row().cells
        row[0].text = member
        row[1].text = role

    add_doc_heading(doc, "14. Conclusion", 1)
    add_doc_paragraph(
        doc,
        "The project successfully demonstrates secure data transmission over a Netcat-style tool. "
        "Normal Netcat sends data as plaintext, and the secret message is visible in Wireshark. "
        "Ncat SSL encrypts the channel, so the same message reaches the receiver but is not readable in packet capture. "
        "Therefore, encrypted channels protect confidentiality during transmission.",
    )

    add_doc_heading(doc, "15. References", 1)
    add_doc_bullets(
        doc,
        [
            "Netcat manual pages",
            "Nmap Ncat documentation: https://nmap.org/ncat/",
            "Wireshark documentation: https://www.wireshark.org/docs/",
            "OpenSSL documentation: https://www.openssl.org/docs/",
        ],
    )

    add_doc_heading(doc, "Appendix: Screenshot Checklist", 1)
    add_doc_bullets(
        doc,
        [
            "Plaintext receiver command screenshot.",
            "Plaintext sender command screenshot.",
            "Wireshark TCP stream showing readable message on port 4444.",
            "Certificate generation screenshot.",
            "Encrypted receiver command screenshot.",
            "Encrypted sender command screenshot.",
            "Wireshark TCP stream showing unreadable data on port 4445.",
        ],
    )

    doc.save(OUT_DIR / "Secure_Data_Transmission_over_Netcat_Report.docx")


def set_slide_title(shape, text):
    shape.text = text
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.size = PptPt(30)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)


def add_bullets(slide, bullets, left=0.8, top=1.55, width=8.7, height=4.8, font_size=22):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(font_size)
    return box


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def create_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = TOPIC
    slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(34)
    slide.placeholders[1].text = f"Practical Exploration of Cryptography / Security Tools\nSelected Tool: Netcat / Ncat\nTeam: {TEAM}"
    add_notes(slide, "Introduce the team and say the project compares unsafe plaintext transfer with encrypted transfer.")

    slides = [
        (
            "AAT Requirement Mapping",
            [
                "Tool introduction: Netcat / Ncat",
                "Setup: Netcat, Ncat, Wireshark, OpenSSL",
                "Practical demo: send same file in two ways",
                "Cryptanalysis: capture traffic using Wireshark",
                "Relevance: confidentiality in network security",
            ],
            "This slide shows that the project follows the assignment instructions.",
        ),
        (
            "Problem Statement",
            [
                "Data sent over a network can be captured by attackers.",
                "If the channel is plaintext, the message is readable.",
                "Confidential data must be protected using encryption.",
                "Our project proves this with a simple Netcat demo.",
            ],
            "Explain the problem in simple language: readable packets are unsafe.",
        ),
        (
            "What is Netcat?",
            [
                "Netcat is a command-line networking tool.",
                "It can send and receive data over TCP/UDP.",
                "It is useful for testing, file transfer, and debugging.",
                "Traditional Netcat does not encrypt data.",
            ],
            "Tanishk can explain this slide.",
        ),
        (
            "Key Security Concept: Confidentiality",
            [
                "Confidentiality means only authorized users can read the data.",
                "Plaintext = readable data.",
                "Ciphertext = encrypted, unreadable data.",
                "Encryption converts plaintext into ciphertext.",
            ],
            "Aastha can explain confidentiality and encryption.",
        ),
        (
            "Tools Used",
            [
                "Netcat (nc): plaintext sender/receiver channel",
                "Ncat SSL: encrypted Netcat-compatible channel",
                "Wireshark: packet capture and traffic inspection",
                "OpenSSL: creates a temporary certificate",
            ],
            "Mention that Ncat is used because classical Netcat has no encryption.",
        ),
        (
            "Mini Project Workflow",
            [
                "Step 1: Create a secret message file.",
                "Step 2: Send it using normal Netcat on port 4444.",
                "Step 3: Capture and read the message in Wireshark.",
                "Step 4: Send same file using Ncat SSL on port 4445.",
                "Step 5: Capture again and compare readability.",
            ],
            "This is the complete practical demonstration in one slide.",
        ),
        (
            "Experiment 1: Plaintext Netcat",
            [
                "Receiver: nc -lvnp 4444 > received_plain.txt",
                "Sender: nc 127.0.0.1 4444 < secret_message.txt",
                "Wireshark filter: tcp.port == 4444",
                "Observation: message is visible in Follow TCP Stream.",
            ],
            "This is the unsafe channel demonstration.",
        ),
        (
            "Cryptanalysis: Packet Sniffing",
            [
                "Attacker opens Wireshark and captures packets.",
                "Attacker follows the TCP stream.",
                "In plaintext mode, the secret message appears directly.",
                "No decryption skill is needed because there is no encryption.",
            ],
            "Aditi can explain that this is the attack/demo of cryptanalysis.",
        ),
        (
            "Experiment 2: Encrypted Ncat SSL",
            [
                "Generate certificate using OpenSSL.",
                "Receiver: ncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445",
                "Sender: ncat --ssl 127.0.0.1 4445 < secret_message.txt",
                "Wireshark filter: tcp.port == 4445",
            ],
            "This is the secure channel demonstration.",
        ),
        (
            "Encrypted Channel Observation",
            [
                "Receiver still gets the correct message.",
                "Wireshark does not show the readable message.",
                "Captured traffic appears as encrypted TLS/SSL data.",
                "Confidentiality is protected.",
            ],
            "Say that delivery works but the attacker cannot read the data.",
        ),
    ]

    for title, bullets, notes in slides:
        slide = prs.slides.add_slide(content_layout)
        set_slide_title(slide.shapes.title, title)
        add_bullets(slide, bullets)
        add_notes(slide, notes)

    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide.shapes.title, "Plaintext vs Encrypted Comparison")
    table = slide.shapes.add_table(6, 3, PptInches(0.7), PptInches(1.45), PptInches(12), PptInches(4.4)).table
    headers = ["Feature", "Plaintext Netcat", "Encrypted Ncat SSL"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    data = [
        ("Port", "4444", "4445"),
        ("Readable in Wireshark", "Yes", "No"),
        ("Confidentiality", "Not protected", "Protected"),
        ("Attacker can read", "Yes", "No"),
        ("Use for sensitive data", "Unsafe", "Safer"),
    ]
    for r, row_values in enumerate(data, start=1):
        for c, value in enumerate(row_values):
            table.cell(r, c).text = value
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.font.size = PptPt(16)
    add_notes(slide, "Use this slide for marks because it clearly shows the result.")

    more_slides = [
        (
            "Applications / Relevance",
            [
                "HTTPS protects website communication.",
                "SSH protects remote login.",
                "VPNs protect traffic on unsafe networks.",
                "Secure file transfer protects files in transit.",
                "Encrypted messaging protects private chats.",
            ],
            "Connect the demo to real-world cybersecurity.",
        ),
        (
            "Limitations",
            [
                "Traditional Netcat does not encrypt by default.",
                "Ncat SSL is mainly used here for learning/demo.",
                "Self-signed certificates are okay for lab demo only.",
                "Production systems should use SSH, SCP, HTTPS, or VPN.",
            ],
            "Be honest about limitations; this improves credibility.",
        ),
        (
            "Conclusion",
            [
                "Normal Netcat sends plaintext data.",
                "Plaintext data is readable in Wireshark.",
                "Ncat SSL encrypts the data channel.",
                "Encrypted traffic protects confidentiality.",
                "Therefore, encryption is necessary for secure transmission.",
            ],
            "End confidently with the key result.",
        ),
        (
            "Team Speaking Split",
            [
                "Tanishk: Netcat introduction and plaintext demo.",
                "Aastha: Confidentiality and encrypted Ncat SSL demo.",
                "Aditi: Wireshark observation, comparison, conclusion.",
                "Everyone: Be ready for basic viva questions.",
            ],
            "This helps the team divide the presentation clearly.",
        ),
    ]

    for title, bullets, notes in more_slides:
        slide = prs.slides.add_slide(content_layout)
        set_slide_title(slide.shapes.title, title)
        add_bullets(slide, bullets)
        add_notes(slide, notes)

    prs.save(OUT_DIR / "Secure_Data_Transmission_over_Netcat_Presentation.pptx")


def write_text_files():
    (OUT_DIR / "REPORT_TEXT.md").write_text(REPORT_MD, encoding="utf-8")
    (OUT_DIR / "DEMO_GUIDE.md").write_text(DEMO_GUIDE, encoding="utf-8")
    (OUT_DIR / "VIVA_QA.md").write_text(VIVA_QA, encoding="utf-8")
    (OUT_DIR / "demo-files" / "secret_message.txt").write_text(
        "This is confidential CRP data from Tanishk, Aastha and Aditi.\n"
        "If this line is visible in Wireshark, confidentiality is broken.\n",
        encoding="utf-8",
    )


def main():
    write_text_files()
    create_report_docx()
    create_pptx()


if __name__ == "__main__":
    main()
