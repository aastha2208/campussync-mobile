
from pathlib import Path
import math
import random

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt
from PIL import Image, ImageDraw, ImageFilter

OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "assets"
ASSET_DIR.mkdir(exist_ok=True)
TEAM = "Tanishk, Aastha, Aditi"
TOPIC = "Secure Data Transmission over Netcat: Demonstrating Confidentiality through Plaintext vs Encrypted Channels"
REPORT_PATH = OUT_DIR / "Secure_Data_Transmission_over_Netcat_Report.docx"
PPT_PATH = OUT_DIR / "Secure_Data_Transmission_over_Netcat_Presentation.pptx"
BG_PATH = ASSET_DIR / "network_security_background.png"


def create_background():
    width, height = 1920, 1080
    img = Image.new("RGB", (width, height), "#07111f")
    pix = img.load()
    for y in range(height):
        for x in range(width):
            nx = x / width
            ny = y / height
            glow = int(42 * (1 - math.sqrt((nx - 0.72) ** 2 + (ny - 0.25) ** 2)))
            r = max(5, 5 + glow // 4)
            g = max(15, 18 + glow // 2)
            b = max(30, 42 + glow)
            pix[x, y] = (min(r, 28), min(g, 84), min(b, 130))

    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    random.seed(42)
    nodes = []
    for _ in range(85):
        x = random.randint(80, width - 80)
        y = random.randint(70, height - 70)
        nodes.append((x, y))

    for i, (x1, y1) in enumerate(nodes):
        for x2, y2 in nodes[i + 1:]:
            dist = math.hypot(x1 - x2, y1 - y2)
            if dist < 230 and random.random() < 0.34:
                alpha = int(max(28, 120 - dist * 0.35))
                draw.line((x1, y1, x2, y2), fill=(38, 196, 255, alpha), width=2)

    for x, y in nodes:
        radius = random.randint(3, 7)
        draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=(117, 222, 255, 210))
        draw.ellipse((x - radius * 3, y - radius * 3, x + radius * 3, y + radius * 3), outline=(117, 222, 255, 60), width=2)

    # Add subtle lock outline.
    lock_x, lock_y = 1450, 610
    draw.rounded_rectangle((lock_x, lock_y, lock_x + 260, lock_y + 210), radius=32, outline=(118, 221, 255, 110), width=7)
    draw.arc((lock_x + 55, lock_y - 130, lock_x + 205, lock_y + 60), 180, 360, fill=(118, 221, 255, 115), width=9)
    draw.line((lock_x + 55, lock_y - 35, lock_x + 55, lock_y + 5), fill=(118, 221, 255, 115), width=9)
    draw.line((lock_x + 205, lock_y - 35, lock_x + 205, lock_y + 5), fill=(118, 221, 255, 115), width=9)
    draw.ellipse((lock_x + 118, lock_y + 80, lock_x + 142, lock_y + 104), fill=(118, 221, 255, 130))
    draw.rounded_rectangle((lock_x + 126, lock_y + 98, lock_x + 134, lock_y + 145), radius=4, fill=(118, 221, 255, 130))

    overlay = overlay.filter(ImageFilter.GaussianBlur(radius=0.25))
    img = Image.alpha_composite(img.convert("RGBA"), overlay)
    vignette = Image.new("L", (width, height), 0)
    vd = ImageDraw.Draw(vignette)
    vd.ellipse((-250, -180, width + 250, height + 180), fill=210)
    vignette = vignette.filter(ImageFilter.GaussianBlur(120))
    dark = Image.new("RGBA", (width, height), (0, 0, 0, 120))
    img = Image.composite(img, Image.alpha_composite(img, dark), vignette.point(lambda v: 255 - v))
    img.convert("RGB").save(BG_PATH, quality=95)


REPORT_MD = f"""# {TOPIC}

**Assignment Title:** Practical Exploration of Cryptography / Security Tools  
**Selected Tool:** Netcat / Ncat  
**Team Members:** {TEAM}

## Abstract
This mini project demonstrates confidentiality in network communication using a simple sender-receiver model. First, a secret message is transmitted using normal Netcat. Wireshark is used to capture the packets, and the message is visible in readable form. This proves that normal Netcat communication is plaintext and not suitable for confidential data. Next, the same message is transmitted using Ncat SSL, a Netcat-compatible encrypted channel. The receiver still receives the correct message, but Wireshark cannot display the original message in readable form. This proves that encryption protects confidentiality during transmission.

## 1. Introduction of the Tool
Netcat is a command-line networking utility used to read and write data across network connections using TCP or UDP. It is often called a networking Swiss Army knife because it can create quick client-server connections, transfer files, test ports, and debug network services.

For this project, Netcat is used because it is simple and clearly shows what happens when data is sent over a network. Traditional Netcat does not provide encryption. Therefore, to demonstrate the encrypted version, we use Ncat with SSL. Ncat is a modern Netcat-compatible tool from the Nmap project and supports SSL/TLS encryption.

## 2. Problem Statement
In real networks, data packets travel through cables, Wi-Fi, routers, and switches. If the communication is not encrypted, a person with packet-capturing access may read the data. This is dangerous for passwords, private messages, login data, confidential documents, and financial information. The project demonstrates this risk using normal Netcat and then shows how an encrypted channel protects the same data.

## 3. Objectives
1. Understand basic Netcat sender-receiver communication.
2. Send a secret text file through a plaintext Netcat channel.
3. Capture plaintext traffic using Wireshark and prove that the message is readable.
4. Send the same file through an encrypted Ncat SSL channel.
5. Capture encrypted traffic and prove that the original message is not readable.
6. Relate the demonstration to confidentiality in cryptography and network security.

## 4. Marking Scheme Alignment
| AAT Requirement | How Our Project Covers It |
|---|---|
| Understanding of tool and usage | Explains Netcat, Ncat SSL, ports, sender, receiver, and commands. |
| Practical demonstration | Shows plaintext and encrypted file transfer. |
| Cryptanalysis demonstration | Uses Wireshark packet sniffing to compare readable and unreadable traffic. |
| Application / relevance | Connects the demo to HTTPS, SSH, VPN, and secure file transfer. |
| Report | Includes theory, setup, methodology, commands, observations, screenshots, and conclusion. |
| Oral presentation | PPT, speaking script, and viva Q&A are prepared. |

## 5. Tools and Setup
| Tool | Purpose |
|---|---|
| Netcat (`nc`) | Creates plaintext sender and receiver connection. |
| Ncat (`ncat --ssl`) | Creates encrypted Netcat-compatible connection. |
| Wireshark | Captures and analyzes packets. |
| OpenSSL | Generates a temporary SSL certificate and key. |
| Two terminals | One terminal acts as receiver and one terminal acts as sender. |

Installation on Ubuntu/Linux:
```bash
sudo apt update
sudo apt install netcat-openbsd ncat wireshark openssl -y
```

Recommended beginner setup: use one laptop with two terminals and Wireshark. Use `127.0.0.1` as the receiver IP address. If Wireshark does not show loopback traffic on Windows, use two laptops on the same Wi-Fi network.

## 6. Basic Theory
**Plaintext** is readable data. Example: `This is confidential CRP data`.  
**Ciphertext** is encrypted unreadable data.  
**Confidentiality** means only authorized users should be able to read the data.  
**Encryption** converts plaintext into ciphertext so that captured packets do not reveal the original message.

Normal Netcat sends data as plaintext. This means Wireshark can show the message using Follow TCP Stream. Ncat SSL encrypts the channel using SSL/TLS, so Wireshark sees network packets but not the original readable message.

## 7. Project Architecture
The project contains three roles:
1. Sender: sends a secret file.
2. Receiver: receives the file.
3. Observer/attacker: captures traffic in Wireshark.

Plaintext flow: Sender -> Netcat plaintext channel -> Receiver. Wireshark can read the message.  
Encrypted flow: Sender -> Ncat SSL encrypted channel -> Receiver. Wireshark cannot read the message.

## 8. Practical Demonstration: Secret File
Create the secret file:
```bash
echo "This is confidential CRP data from Tanishk, Aastha and Aditi." > secret_message.txt
echo "If this line is visible in Wireshark, confidentiality is broken." >> secret_message.txt
```

## 9. Experiment 1: Plaintext Netcat Channel
Start Wireshark and use filter:
```text
tcp.port == 4444
```
Receiver terminal:
```bash
nc -lvnp 4444 > received_plain.txt
```
Sender terminal:
```bash
nc 127.0.0.1 4444 < secret_message.txt
```
Check received file:
```bash
cat received_plain.txt
```
Wireshark observation: right-click a packet, select Follow > TCP Stream. The secret message is visible in readable form.

## 10. Cryptanalysis Demonstration
For this project, cryptanalysis is shown through packet sniffing. In the plaintext channel, there is no encryption to break. The observer simply captures packets in Wireshark, follows the TCP stream, and reads the message. This proves that plaintext communication fails to provide confidentiality.

## 11. Experiment 2: Encrypted Ncat SSL Channel
Generate a temporary SSL certificate:
```bash
openssl req -x509 -newkey rsa:2048 -nodes -keyout server.key -out server.crt -days 1 -subj "/CN=NetcatAAT"
```
Start Wireshark and use filter:
```text
tcp.port == 4445
```
Encrypted receiver terminal:
```bash
ncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445 > received_secure.txt
```
Encrypted sender terminal:
```bash
ncat --ssl 127.0.0.1 4445 < secret_message.txt
```
Check received file:
```bash
cat received_secure.txt
```
Wireshark observation: Follow TCP Stream does not show the readable secret message. The captured data appears encrypted.

## 12. Results and Comparison
| Feature | Plaintext Netcat | Encrypted Ncat SSL |
|---|---|---|
| Port used | 4444 | 4445 |
| Tool command | `nc` | `ncat --ssl` |
| Receiver gets message | Yes | Yes |
| Wireshark can read message | Yes | No |
| Confidentiality | Not protected | Protected |
| Attacker effort | Very easy | Cannot read without decryption key |
| Suitable for sensitive data | No | Better for demonstration of secure channel |

## 13. Application and Relevance
This project relates directly to real-world security. HTTPS encrypts website traffic, SSH encrypts remote login, VPNs encrypt network traffic on unsafe networks, and secure file transfer tools encrypt files during transmission. The same basic idea is demonstrated in this project: readable plaintext is unsafe, while encrypted data protects confidentiality.

## 14. Limitations
1. Traditional Netcat does not support encryption by default.
2. Ncat SSL is used for learning and demonstration, not as a replacement for well-managed production systems.
3. Self-signed certificates are acceptable for lab demos but not recommended for public production services.
4. This project mainly focuses on confidentiality; authentication and integrity can be discussed as future improvements.

## 15. Future Scope
The project can be extended by adding certificate verification, comparing SSH/SCP with Netcat, transferring larger files, calculating hashes before and after transfer to prove integrity, or demonstrating how HTTPS uses TLS similarly.

## 16. Team Contribution
| Member | Responsibility |
|---|---|
| Tanishk | Netcat introduction, setup, plaintext demo. |
| Aastha | Confidentiality theory, encryption explanation, Ncat SSL demo. |
| Aditi | Wireshark analysis, comparison table, report/PPT conclusion. |

## 17. Conclusion
The project successfully demonstrates the importance of encryption in secure data transmission. Normal Netcat sends the message as plaintext, and Wireshark can read it easily. Ncat SSL sends the same message through an encrypted channel, so Wireshark cannot show the readable message. Therefore, encrypted channels protect confidentiality and are necessary when sensitive data is transmitted over a network.

## References
1. Nmap Ncat documentation: https://nmap.org/ncat/
2. Wireshark documentation: https://www.wireshark.org/docs/
3. OpenSSL documentation: https://www.openssl.org/docs/
4. Netcat manual pages

## Screenshot Checklist
1. Plaintext receiver command.
2. Plaintext sender command.
3. Wireshark Follow TCP Stream showing readable message on port 4444.
4. OpenSSL certificate generation command.
5. Encrypted receiver command.
6. Encrypted sender command.
7. Wireshark Follow TCP Stream showing unreadable encrypted data on port 4445.
"""


def add_run(paragraph, text, bold=False, size=None, color=None):
    run = paragraph.add_run(text)
    run.bold = bold
    if size:
        run.font.size = Pt(size)
    if color:
        run.font.color.rgb = RGBColor(*color)
    return run


def add_code(doc, code):
    for line in code.strip("\n").splitlines():
        p = doc.add_paragraph()
        r = p.add_run(line)
        r.font.name = "Courier New"
        r.font.size = Pt(9)


def add_bullets(doc, items):
    for item in items:
        doc.add_paragraph(item, style="List Bullet")


def create_report():
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.7)
    sec.bottom_margin = Inches(0.7)
    sec.left_margin = Inches(0.8)
    sec.right_margin = Inches(0.8)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_run(p, "ALTERNATE ASSESSMENT TOOL REPORT", True, 16, (31, 78, 121))
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_run(p, TOPIC, True, 14)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_run(p, f"Assignment Title: Practical Exploration of Cryptography / Security Tools\nSelected Tool: Netcat / Ncat\nTeam Members: {TEAM}", True, 11)

    sections = [
        ("Abstract", ["This mini project demonstrates confidentiality in network communication using a simple sender-receiver model. First, a secret message is transmitted using normal Netcat. Wireshark is used to capture the packets, and the message is visible in readable form. This proves that normal Netcat communication is plaintext and not suitable for confidential data. Next, the same message is transmitted using Ncat SSL, a Netcat-compatible encrypted channel. The receiver still receives the correct message, but Wireshark cannot display the original message in readable form. This proves that encryption protects confidentiality during transmission."]),
        ("1. Introduction of the Tool", ["Netcat is a command-line networking utility used to read and write data across network connections using TCP or UDP. It can create quick client-server connections, transfer files, test ports, and debug simple network services.", "Traditional Netcat does not provide encryption. Therefore, this project uses Ncat with SSL for the encrypted part. Ncat is a modern Netcat-compatible tool from the Nmap project and supports SSL/TLS encryption."]),
        ("2. Problem Statement", ["In real networks, data packets travel through Wi-Fi, routers, and switches. If communication is not encrypted, a person with packet-capturing access may read the data. This is dangerous for passwords, private messages, login data, confidential documents, and financial information."]),
        ("3. Objectives", []),
    ]
    for title, paragraphs in sections:
        doc.add_heading(title, level=1)
        for text in paragraphs:
            doc.add_paragraph(text)
    add_bullets(doc, ["Understand basic Netcat sender-receiver communication.", "Send a secret text file through a plaintext Netcat channel.", "Capture plaintext traffic using Wireshark and prove that the message is readable.", "Send the same file through an encrypted Ncat SSL channel.", "Capture encrypted traffic and prove that the original message is not readable.", "Relate the demonstration to confidentiality in cryptography and network security."])

    doc.add_heading("4. Marking Scheme Alignment", level=1)
    t = doc.add_table(rows=1, cols=2)
    t.style = "Table Grid"
    t.rows[0].cells[0].text = "AAT Requirement"
    t.rows[0].cells[1].text = "How Our Project Covers It"
    for a, b in [("Understanding of tool and usage", "Explains Netcat, Ncat SSL, ports, sender, receiver, and commands."), ("Practical demonstration", "Shows plaintext and encrypted file transfer."), ("Cryptanalysis demonstration", "Uses Wireshark packet sniffing to compare readable and unreadable traffic."), ("Application / relevance", "Connects the demo to HTTPS, SSH, VPN, and secure file transfer."), ("Report", "Includes theory, setup, methodology, commands, observations, screenshots, and conclusion."), ("Oral presentation", "PPT, speaking script, and viva Q&A are prepared.")]:
        row = t.add_row().cells
        row[0].text = a
        row[1].text = b

    doc.add_heading("5. Tools and Setup", level=1)
    t = doc.add_table(rows=1, cols=2)
    t.style = "Table Grid"
    t.rows[0].cells[0].text = "Tool"
    t.rows[0].cells[1].text = "Purpose"
    for a, b in [("Netcat (nc)", "Creates plaintext sender and receiver connection."), ("Ncat (ncat --ssl)", "Creates encrypted Netcat-compatible connection."), ("Wireshark", "Captures and analyzes packets."), ("OpenSSL", "Generates a temporary SSL certificate and key."), ("Two terminals", "One terminal acts as receiver and one terminal acts as sender.")]:
        row = t.add_row().cells
        row[0].text = a
        row[1].text = b
    doc.add_paragraph("Ubuntu/Linux installation command:")
    add_code(doc, "sudo apt update\nsudo apt install netcat-openbsd ncat wireshark openssl -y")
    doc.add_paragraph("Recommended beginner setup: one laptop with two terminals and Wireshark. Use 127.0.0.1 as the receiver IP address. If loopback capture is difficult on Windows, use two laptops on the same Wi-Fi network.")

    doc.add_heading("6. Basic Theory", level=1)
    add_bullets(doc, ["Plaintext: readable data before encryption.", "Ciphertext: encrypted unreadable data.", "Confidentiality: only authorized users should be able to read the data.", "Encryption: converts plaintext into ciphertext so captured packets do not reveal the original message."])
    doc.add_paragraph("Normal Netcat sends data as plaintext. Ncat SSL encrypts the channel using SSL/TLS, so Wireshark sees network packets but not the original readable message.")

    doc.add_heading("7. Project Architecture", level=1)
    add_bullets(doc, ["Sender: sends a secret file.", "Receiver: receives the file.", "Observer/attacker: captures traffic in Wireshark."])
    doc.add_paragraph("Plaintext flow: Sender -> Netcat plaintext channel -> Receiver. Wireshark can read the message.")
    doc.add_paragraph("Encrypted flow: Sender -> Ncat SSL encrypted channel -> Receiver. Wireshark cannot read the message.")

    doc.add_heading("8. Practical Demonstration: Secret File", level=1)
    add_code(doc, 'echo "This is confidential CRP data from Tanishk, Aastha and Aditi." > secret_message.txt\necho "If this line is visible in Wireshark, confidentiality is broken." >> secret_message.txt')

    doc.add_heading("9. Experiment 1: Plaintext Netcat Channel", level=1)
    doc.add_paragraph("Start Wireshark and use display filter: tcp.port == 4444")
    doc.add_paragraph("Receiver terminal:")
    add_code(doc, "nc -lvnp 4444 > received_plain.txt")
    doc.add_paragraph("Sender terminal:")
    add_code(doc, "nc 127.0.0.1 4444 < secret_message.txt")
    doc.add_paragraph("Observation: In Wireshark, Follow TCP Stream shows the secret message in readable form. This proves confidentiality is broken.")

    doc.add_heading("10. Cryptanalysis Demonstration", level=1)
    doc.add_paragraph("In this project, cryptanalysis is shown through packet sniffing. In the plaintext channel, there is no encryption to break. The observer simply captures packets in Wireshark, follows the TCP stream, and reads the message.")

    doc.add_heading("11. Experiment 2: Encrypted Ncat SSL Channel", level=1)
    doc.add_paragraph("Generate temporary SSL certificate:")
    add_code(doc, 'openssl req -x509 -newkey rsa:2048 -nodes -keyout server.key -out server.crt -days 1 -subj "/CN=NetcatAAT"')
    doc.add_paragraph("Start Wireshark and use display filter: tcp.port == 4445")
    doc.add_paragraph("Encrypted receiver terminal:")
    add_code(doc, "ncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445 > received_secure.txt")
    doc.add_paragraph("Encrypted sender terminal:")
    add_code(doc, "ncat --ssl 127.0.0.1 4445 < secret_message.txt")
    doc.add_paragraph("Observation: The receiver gets the correct message, but Wireshark does not show the readable secret text. The traffic appears encrypted.")

    doc.add_heading("12. Results and Comparison", level=1)
    t = doc.add_table(rows=1, cols=3)
    t.style = "Table Grid"
    for i, h in enumerate(["Feature", "Plaintext Netcat", "Encrypted Ncat SSL"]):
        t.rows[0].cells[i].text = h
    for vals in [("Port used", "4444", "4445"), ("Tool command", "nc", "ncat --ssl"), ("Receiver gets message", "Yes", "Yes"), ("Wireshark can read message", "Yes", "No"), ("Confidentiality", "Not protected", "Protected"), ("Attacker effort", "Very easy", "Cannot read without decryption key"), ("Suitable for sensitive data", "No", "Better for secure-channel demonstration")]:
        row = t.add_row().cells
        for i, val in enumerate(vals):
            row[i].text = val

    doc.add_heading("13. Application and Relevance", level=1)
    doc.add_paragraph("This project relates directly to real-world security. HTTPS encrypts website traffic, SSH encrypts remote login, VPNs encrypt network traffic on unsafe networks, and secure file transfer tools encrypt files during transmission. The same basic idea is demonstrated here: readable plaintext is unsafe, while encrypted data protects confidentiality.")

    doc.add_heading("14. Limitations", level=1)
    add_bullets(doc, ["Traditional Netcat does not support encryption by default.", "Ncat SSL is used for learning and demonstration, not as a replacement for well-managed production systems.", "Self-signed certificates are acceptable for lab demos but not recommended for public production services.", "This project mainly focuses on confidentiality; authentication and integrity can be future improvements."])

    doc.add_heading("15. Future Scope", level=1)
    add_bullets(doc, ["Add certificate verification.", "Compare SSH/SCP with Netcat.", "Transfer larger files and compare performance.", "Use file hashes before and after transfer to prove integrity.", "Explain how HTTPS uses TLS in a similar way."])

    doc.add_heading("16. Team Contribution", level=1)
    t = doc.add_table(rows=1, cols=2)
    t.style = "Table Grid"
    t.rows[0].cells[0].text = "Member"
    t.rows[0].cells[1].text = "Responsibility"
    for a, b in [("Tanishk", "Netcat introduction, setup, plaintext demo."), ("Aastha", "Confidentiality theory, encryption explanation, Ncat SSL demo."), ("Aditi", "Wireshark analysis, comparison table, report/PPT conclusion.")]:
        row = t.add_row().cells
        row[0].text = a
        row[1].text = b

    doc.add_heading("17. Conclusion", level=1)
    doc.add_paragraph("The project successfully demonstrates the importance of encryption in secure data transmission. Normal Netcat sends the message as plaintext, and Wireshark can read it easily. Ncat SSL sends the same message through an encrypted channel, so Wireshark cannot show the readable message. Therefore, encrypted channels protect confidentiality and are necessary when sensitive data is transmitted over a network.")

    doc.add_heading("References", level=1)
    add_bullets(doc, ["Nmap Ncat documentation: https://nmap.org/ncat/", "Wireshark documentation: https://www.wireshark.org/docs/", "OpenSSL documentation: https://www.openssl.org/docs/", "Netcat manual pages"])

    doc.add_heading("Screenshot Checklist", level=1)
    add_bullets(doc, ["Plaintext receiver command.", "Plaintext sender command.", "Wireshark Follow TCP Stream showing readable message on port 4444.", "OpenSSL certificate generation command.", "Encrypted receiver command.", "Encrypted sender command.", "Wireshark Follow TCP Stream showing unreadable encrypted data on port 4445."])
    doc.save(REPORT_PATH)


def add_bg(slide, prs):
    slide.shapes.add_picture(str(BG_PATH), 0, 0, width=prs.slide_width, height=prs.slide_height)


def title_box(slide, title):
    box = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.33), PptInches(12.2), PptInches(0.72))
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = PptRGBColor(4, 15, 30)
    fill.transparency = 18
    line = box.line
    line.color.rgb = PptRGBColor(75, 205, 255)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = PptPt(25)
    p.font.color.rgb = PptRGBColor(232, 248, 255)
    return box


def content_box(slide, bullets, left=0.75, top=1.35, width=11.85, height=5.55, font_size=20):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = PptRGBColor(5, 18, 34)
    fill.transparency = 12
    box.line.color.rgb = PptRGBColor(34, 130, 180)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = PptInches(0.25)
    tf.margin_right = PptInches(0.22)
    tf.margin_top = PptInches(0.18)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = PptPt(6)
        p.font.size = PptPt(font_size)
        p.font.color.rgb = PptRGBColor(236, 250, 255)
    return box


def code_box(slide, title, code, left=0.85, top=2.1, width=11.65, height=3.7):
    content = [title] + code.strip("\n").splitlines()
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = PptRGBColor(0, 8, 16)
    fill.transparency = 5
    box.line.color.rgb = PptRGBColor(89, 221, 255)
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Courier New" if i else "Aptos"
        p.font.bold = i == 0
        p.font.size = PptPt(18 if i else 21)
        p.font.color.rgb = PptRGBColor(160, 242, 255) if i else PptRGBColor(255, 255, 255)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def create_ppt():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs)
    title = slide.shapes.add_textbox(PptInches(0.7), PptInches(1.55), PptInches(12.0), PptInches(1.65))
    title.fill.solid(); title.fill.fore_color.rgb = PptRGBColor(3, 16, 31); title.fill.transparency = 8
    title.line.color.rgb = PptRGBColor(87, 215, 255)
    tf = title.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = TOPIC; p.alignment = PP_ALIGN.CENTER; p.font.bold = True; p.font.size = PptPt(31); p.font.color.rgb = PptRGBColor(236, 250, 255)
    sub = slide.shapes.add_textbox(PptInches(1.7), PptInches(3.55), PptInches(10.0), PptInches(1.0))
    sub.fill.solid(); sub.fill.fore_color.rgb = PptRGBColor(3, 16, 31); sub.fill.transparency = 18
    stf = sub.text_frame; stf.clear()
    p = stf.paragraphs[0]; p.text = f"Practical Exploration of Cryptography / Security Tools\nSelected Tool: Netcat / Ncat | Team: {TEAM}"; p.alignment = PP_ALIGN.CENTER; p.font.size = PptPt(21); p.font.color.rgb = PptRGBColor(220, 244, 255)
    notes(slide, "Introduce the team and say the project compares plaintext communication with encrypted communication.")

    slides = [
        ("Why We Chose This Topic", ["Easy to demonstrate using simple commands.", "Directly connected to confidentiality in cryptography.", "Shows a real security problem: packet sniffing.", "Gives clear visual proof through Wireshark screenshots.", "Suitable for a 10-mark practical AAT because it includes tool usage, demo, analysis, and conclusion."], "Explain that the topic is practical and understandable."),
        ("AAT Requirement Mapping", ["Introduction of tool: Netcat/Ncat explained.", "Installation/setup: commands and required software included.", "Practical demonstration: plaintext vs encrypted transfer.", "Cryptanalysis: Wireshark packet sniffing and TCP stream inspection.", "Application/relevance: HTTPS, SSH, VPN, secure file transfer.", "Report + oral presentation: ready report, PPT, script, viva Q&A."], "This slide directly maps to the professor's rules."),
        ("Problem Statement", ["Data sent over a network travels as packets.", "If packets are not encrypted, captured traffic may reveal the message.", "Attackers can use packet sniffers on unsafe networks.", "We demonstrate this risk with normal Netcat.", "Then we demonstrate confidentiality using encrypted Ncat SSL."], "Explain the real-world need for encryption."),
        ("What is Netcat?", ["Netcat is a command-line networking utility.", "It can create simple TCP/UDP sender-receiver connections.", "It is used for testing ports, sending messages, and transferring files.", "Normal Netcat sends data as plaintext.", "Because it is simple, it is perfect for learning network security basics."], "Tanishk can present this."),
        ("Why Use Ncat SSL?", ["Traditional Netcat does not encrypt data.", "Ncat is a Netcat-compatible tool from the Nmap project.", "Ncat supports SSL/TLS using the --ssl option.", "This lets us compare an unsafe channel and an encrypted channel clearly.", "The project still stays within the Netcat-family tool idea."], "Clarify why Ncat appears in the encrypted part."),
        ("Security Concept: Confidentiality", ["Confidentiality means only authorized users can read the data.", "Plaintext is readable data before encryption.", "Ciphertext is encrypted, unreadable data.", "Encryption converts plaintext into ciphertext.", "Our test: if Wireshark can read the secret message, confidentiality failed."], "Aastha can present the security theory."),
        ("Tools Used", ["Netcat (nc): plaintext transfer on port 4444.", "Ncat SSL: encrypted transfer on port 4445.", "Wireshark: captures and inspects packets.", "OpenSSL: generates a temporary SSL certificate.", "Two terminals: one receiver and one sender."], "List tools and their roles."),
        ("Mini Project Architecture", ["Sender: sends secret_message.txt.", "Receiver: saves the received file.", "Observer/attacker: captures packets in Wireshark.", "Plaintext path: Sender -> nc -> Receiver; Wireshark reads message.", "Encrypted path: Sender -> ncat --ssl -> Receiver; Wireshark cannot read message."], "Explain the three roles."),
        ("Demo Setup", ["Use one laptop with two terminals for easiest demo.", "Use 127.0.0.1 when sender and receiver are on same laptop.", "Use port 4444 for plaintext Netcat.", "Use port 4445 for encrypted Ncat SSL.", "Open Wireshark before sending data to capture packets."], "Make the setup easy for beginners."),
        ("Create Secret Message File", ["The same file is used in both experiments.", "This keeps the comparison fair.", "If this text appears in Wireshark during plaintext transfer, confidentiality is broken.", "If it does not appear during encrypted transfer, confidentiality is protected."], "Prepare before showing commands."),
        ("Plaintext Experiment Commands", ["Receiver: nc -lvnp 4444 > received_plain.txt", "Sender: nc 127.0.0.1 4444 < secret_message.txt", "Wireshark filter: tcp.port == 4444", "Right-click packet -> Follow -> TCP Stream", "Expected result: secret message is visible."], "Show the unsafe experiment."),
        ("Plaintext Observation", ["Receiver gets the file successfully.", "Wireshark also displays the secret text.", "This means an observer can read the data.", "No password or key is needed by the attacker.", "Conclusion: normal Netcat does not provide confidentiality."], "Aditi can explain Wireshark result."),
        ("Cryptanalysis Demonstration", ["We use packet sniffing as the attack method.", "The attacker captures traffic using Wireshark.", "In plaintext mode, there is no encryption to break.", "Follow TCP Stream directly reveals the message.", "This is simple but powerful proof of insecurity."], "This matches the cryptanalysis marks."),
        ("Generate SSL Certificate", ["OpenSSL creates a temporary certificate and private key.", "server.crt is the certificate.", "server.key is the private key.", "For a classroom demo, a self-signed certificate is acceptable.", "The certificate enables the encrypted Ncat SSL listener."], "Explain the certificate in simple terms."),
        ("Encrypted Experiment Commands", ["Receiver: ncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445 > received_secure.txt", "Sender: ncat --ssl 127.0.0.1 4445 < secret_message.txt", "Wireshark filter: tcp.port == 4445", "Right-click packet -> Follow -> TCP Stream", "Expected result: readable message is not visible."], "Show secure experiment."),
        ("Encrypted Observation", ["Receiver still gets the correct file.", "Wireshark captures packets but cannot show the readable message.", "The data appears as encrypted SSL/TLS traffic.", "An observer cannot understand the secret data.", "Conclusion: encryption provides confidentiality."], "State the secure result clearly."),
        ("Plaintext vs Encrypted Result", ["Plaintext Netcat: easy to send, but unsafe for secrets.", "Encrypted Ncat SSL: same message delivered, but hidden from Wireshark.", "Both experiments prove the difference visually.", "This is the main output of our mini project.", "Screenshots should be added here after running the demo."], "Use this slide with screenshots if possible."),
        ("Comparison Table", ["Port: 4444 plaintext, 4445 encrypted.", "Wireshark readability: visible in plaintext, unreadable in encrypted channel.", "Confidentiality: not protected vs protected.", "Attacker result: can read message vs cannot read message.", "Security lesson: sensitive data needs encryption."], "Summarize comparison."),
        ("Applications / Relevance", ["HTTPS protects web browsing and login forms.", "SSH protects remote terminal sessions.", "VPNs protect traffic on unsafe networks.", "Secure file transfer protects documents in transit.", "Encrypted messaging protects private communication."], "Connect to real-world examples."),
        ("Limitations", ["Traditional Netcat has no encryption by default.", "Ncat SSL is good for learning and demonstration.", "Self-signed certificates are not ideal for production.", "Production systems usually use SSH, SCP, HTTPS, or VPN.", "This project mainly focuses on confidentiality."], "Mention limitations honestly."),
        ("Future Scope", ["Add certificate verification.", "Compare the same demo with SSH/SCP.", "Transfer a larger file and compare packet captures.", "Use hashes to verify file integrity.", "Explain how TLS in HTTPS follows a similar concept."], "Future scope helps report quality."),
        ("Team Contribution", ["Tanishk: Netcat introduction, setup, plaintext demo.", "Aastha: confidentiality theory, encryption explanation, Ncat SSL demo.", "Aditi: Wireshark analysis, comparison table, conclusion.", "All members: practice viva answers and screenshots.", "Keep explanation simple and confidence-based."], "Divide the speaking roles."),
        ("Final Conclusion", ["Normal Netcat sends data in plaintext.", "Plaintext data is readable in Wireshark.", "Ncat SSL encrypts the communication channel.", "Encrypted traffic hides the original message from observers.", "Therefore, encryption is necessary for confidential data transmission."], "End with this conclusion."),
    ]

    for title, bullets, note in slides:
        slide = prs.slides.add_slide(blank)
        add_bg(slide, prs)
        title_box(slide, title)
        content_box(slide, bullets)
        notes(slide, note)

    # Add a command-reference slide with monospaced commands.
    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs)
    title_box(slide, "Command Reference for Demo")
    code_box(slide, "Copy these during demo:", 'echo "This is confidential CRP data from Tanishk, Aastha and Aditi." > secret_message.txt\nnc -lvnp 4444 > received_plain.txt\nnc 127.0.0.1 4444 < secret_message.txt\nopenssl req -x509 -newkey rsa:2048 -nodes -keyout server.key -out server.crt -days 1 -subj "/CN=NetcatAAT"\nncat --ssl --ssl-cert server.crt --ssl-key server.key -lvnp 4445 > received_secure.txt\nncat --ssl 127.0.0.1 4445 < secret_message.txt')
    notes(slide, "Keep this as backup if commands are needed during presentation.")

    prs.save(PPT_PATH)


def main():
    create_background()
    (OUT_DIR / "REPORT_TEXT.md").write_text(REPORT_MD, encoding="utf-8")
    create_report()
    create_ppt()


if __name__ == "__main__":
    main()
